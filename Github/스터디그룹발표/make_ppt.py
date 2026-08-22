from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ===== 색상 정의 =====
DARK_NAVY = RGBColor(0x0D, 0x11, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xA4, 0x4F)
GRAY = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_background(slide, color=DARK_NAVY):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 맨 뒤로 보내기
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_title(slide, text, top=Inches(0.5), size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return box


def add_subtitle(slide, text, top, size=20, color=GRAY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return box


def add_bullets(slide, items, top=Inches(2.0), left=Inches(0.9), width=Inches(11.5),
                 size=20, color=WHITE, line_space=1.3):
    box = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(14)
        p.line_spacing = line_space
    return box


def add_note(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_footer_tag(slide, text="실습 자료 별도 배포", color=GREEN):
    box = slide.shapes.add_textbox(Inches(9.5), Inches(0.2), Inches(3.3), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.color.rgb = color
    run.font.bold = True


def add_arrow_flow(slide, labels, top=Inches(3.0), color=GREEN):
    """4단계 프로세스를 화살표 도형으로 표시"""
    n = len(labels)
    box_w = Inches(2.6)
    gap = Inches(0.35)
    total_w = box_w * n + gap * (n - 1)
    start_x = (prs.slide_width - total_w) / 2
    y = top
    h = Inches(1.2)

    for i, label in enumerate(labels):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON if i > 0 else MSO_SHAPE.PENTAGON,
                                        x, y, box_w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Malgun Gothic"


def add_table(slide, data, top=Inches(2.2), left=Inches(1.2), width=Inches(11), height=Inches(3)):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.size = Pt(16)
            run.font.name = "Malgun Gothic"
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREEN
            else:
                run.font.color.rgb = DARK_NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table


# =====================================================
# Slide 1 — 표지
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)

title_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.3), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Git & GitHub, 협업의 시작"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Malgun Gothic"

sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(0.8))
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "마이스터고 학생을 위한 실무 협업 도구 이해하기"
run2.font.size = Pt(22)
run2.font.color.rgb = GREEN
run2.font.name = "Malgun Gothic"

# 하단 라인 포인트 장식
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.9), Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.3), Inches(0.5))
p3 = footer_box.text_frame.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "발표자명 · 2025"
run3.font.size = Pt(14)
run3.font.color.rgb = GRAY

# =====================================================
# Slide 2 — 오늘의 목표
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "오늘 이야기할 것")
add_bullets(slide, [
    "Git이 무엇이고 왜 쓰는가",
    "GitHub로 어떻게 협업하는가",
    "취업 · 포트폴리오와의 연결",
])
add_note(slide, "짧게 3초 내 훑고 바로 다음 슬라이드로 (시간 아끼기)")

# =====================================================
# Slide 3 — 도입: 문제 제기
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "'최종.zip', '진짜최종.zip'... 익숙하시죠?")
add_bullets(slide, [
    "파일명으로 버전 관리하던 경험",
    "누가 뭘 고쳤는지 모르는 혼란",
    "팀 작업 시 파일 덮어쓰기 사고",
])
add_note(slide, "여러분, 최종.zip, 진짜최종.zip 만들어본 적 있죠? → 웃음 유도 후 문제 제기")

# =====================================================
# Slide 4 — Git vs GitHub
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "Git은 일지 기록, GitHub는 일지 공유")
table_data = [
    ["구분", "Git", "GitHub"],
    ["역할", "변경 이력을 기록하는 도구", "그 기록을 온라인에 공유하는 공간"],
    ["비유", "개인 다이어리", "공유 클라우드 다이어리"],
    ["위치", "내 컴퓨터 (로컬)", "인터넷 (원격 서버)"],
]
add_table(slide, table_data, top=Inches(2.3), left=Inches(1.0), width=Inches(11.3), height=Inches(3.5))
add_note(slide, "Git은 이 문제를 자동으로 기록해주는 도구, GitHub는 그걸 공유하는 곳입니다.")

# =====================================================
# Slide 5 — 워크플로우 개요
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "실제 개발자들은 이렇게 일합니다")
add_arrow_flow(slide, ["Commit", "Branch", "Merge", "Pull Request"], top=Inches(3.2))
add_note(slide, "이 네 단계만 이해하면 협업의 90%는 이해한 겁니다.")

# =====================================================
# Slide 6 — Commit
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "Commit = 작업 스냅샷 저장")
add_bullets(slide, [
    "파일을 수정할 때마다 '찰칵' 사진 찍듯 기록",
    "언제, 누가, 무엇을 바꿨는지 로그로 남음",
    "문제 생기면 이전 시점으로 되돌리기 가능",
])

# =====================================================
# Slide 7 — Branch
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "Branch = 나만의 작업 공간")
add_bullets(slide, [
    "원본(main)을 건드리지 않고 복사본에서 작업",
    "여러 사람이 동시에 각자 기능 개발 가능",
    "나무 가지처럼 뻗어나가는 구조",
])

# =====================================================
# Slide 8 — Merge
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "Merge = 작업 결과 합치기")
add_bullets(slide, [
    "완성된 branch를 다시 main에 합침",
    "여러 사람의 작업이 하나로 통합되는 과정",
])

# =====================================================
# Slide 9 — Pull Request
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "PR = '제 작업 확인해주세요' 요청")
add_bullets(slide, [
    "바로 합치지 않고 동료의 검토(리뷰)를 거침",
    "코드 품질 관리 + 협업 커뮤니케이션의 핵심",
    "실무에서 가장 많이 쓰는 협업 방식",
])
add_note(slide, "Commit(저장) → Branch(따로 작업) → Merge(합치기) → PR(검토 요청). 실제 개발자들은 이렇게 일합니다.")

# =====================================================
# Slide 10 — Conflict
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "충돌은 고장이 아니라 신호입니다")
add_bullets(slide, [
    "두 사람이 같은 부분을 동시에 수정하면 발생",
    "Git이 '둘 다 확인해주세요'라고 알려주는 것",
    "실습 화면 예시로만 간단히 보여주기",
])
add_note(slide, "두 사람이 같은 파일을 고치면 충돌이 나요. 근데 이건 고장이 아니라 '둘 다 열심히 했다'는 신호예요.")

# =====================================================
# Slide 11 — 실습 자료 안내
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "실습은 배포 자료를 참고하세요")
add_bullets(slide, [
    "오늘은 개념 이해 중심으로 진행",
    "상세 실습 가이드는 별도 자료로 배포",
])
add_footer_tag(slide)

# =====================================================
# Slide 12 — 포트폴리오 연결 ⭐
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "왜 GitHub이 취업에 도움이 될까?")
add_bullets(slide, [
    "GitHub 프로필 = 실무형 이력서",
    "협업 능력 = 코드 실력만큼 중요한 평가 기준",
    "기업 채용 시 'GitHub 링크'가 차별화 포인트",
])
# 강조 배지
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.6), Inches(1.8), Inches(0.6))
badge.fill.solid()
badge.fill.fore_color.rgb = GREEN
badge.line.fill.background()
badge.text_frame.text = "⭐ 핵심"
badge.text_frame.paragraphs[0].font.size = Pt(16)
badge.text_frame.paragraphs[0].font.bold = True
badge.text_frame.paragraphs[0].font.color.rgb = WHITE
badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_note(slide, "이걸 잘하는 사람이 실무에서 바로 협업 가능한 사람으로 평가받습니다. 포트폴리오에 GitHub 링크 하나 있는 것만으로도 차별화됩니다.")

# =====================================================
# Slide 13 — 실제 예시
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "이런 프로필이 눈에 띕니다")
add_bullets(slide, [
    "꾸준한 커밋 잔디밭 (Contribution Graph)",
    "README로 프로젝트 설명 잘 정리",
    "PR / 이슈에 남긴 협업 흔적",
])

# =====================================================
# Slide 14 — 핵심 요약
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "3줄 요약")
add_bullets(slide, [
    "Git = 기록, GitHub = 공유",
    "Commit → Branch → Merge → PR이 협업의 기본 흐름",
    "GitHub 활동 = 취업 시 강력한 포트폴리오",
], size=24)

# =====================================================
# Slide 15 — Q&A
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "질문 받겠습니다 🙋"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Malgun Gothic"
add_note(slide, "시간 부족 시 '추가 질문은 쉬는 시간에' 안내")

# =====================================================
# Slide 16 — 부록 표지
# =====================================================
slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_title(slide, "📚 부록: 실습 가이드 & 체크리스트")
add_bullets(slide, [
    "실습 슬라이드 (Branch / Merge / Conflict 직접 해보기)",
    "Git 명령어 치트시트",
    "참고 링크 모음",
])
add_note(slide, "발표 중 노출 안 함, 자료 배포용 PDF/PPT 마지막에만 첨부")

# =====================================================
# 저장
# =====================================================
prs.save("Git_GitHub_협업의_시작.pptx")
print("완료: Git_GitHub_협업의_시작.pptx 파일이 생성되었습니다.")